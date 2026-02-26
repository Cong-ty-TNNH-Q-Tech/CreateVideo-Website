"""
Module để đọc nội dung từ file PowerPoint và PDF
"""
from pptx import Presentation
from pypdf import PdfReader
import fitz  # pymupdf
from PIL import Image
import os
import sys
import platform
import subprocess
from typing import List, Dict

class PresentationReader:
    """Đọc nội dung từ file PPT và PDF"""
    
    @staticmethod
    def read_pptx(file_path: str) -> List[Dict]:
        """
        Đọc file PowerPoint và trả về danh sách slides
        """
        try:
            prs = Presentation(file_path)
            slides_data = []
            
            for slide in prs.slides:
                # Check if slide is hidden
                # The 'show' attribute on the slide element indicates visibility (0 = hidden)
                if slide.element.get('show') == '0':
                    continue

                idx = len(slides_data) + 1

                content_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text.strip())
                
                notes = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes = notes_slide.notes_text_frame.text
                
                content = "\n".join(content_parts)
                slides_data.append({
                    'slide_num': idx,
                    'content': content,
                    'notes': notes,
                    'total_slides': len(prs.slides)
                })
            return slides_data
        except Exception as e:
            raise Exception(f"Error reading PPTX: {str(e)}")
    
    @staticmethod
    def read_pdf(file_path: str) -> List[Dict]:
        """
        Đọc file PDF và trả về danh sách pages
        """
        try:
            reader = PdfReader(file_path)
            pages_data = []
            for idx, page in enumerate(reader.pages, 1):
                content = page.extract_text()
                pages_data.append({
                    'slide_num': idx,
                    'content': content.strip(),
                    'notes': "",
                    'total_slides': len(reader.pages)
                })
            return pages_data
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
    
    @staticmethod
    def extract_text_from_file(file_path: str) -> List[Dict]:
        """
        Tự động detect loại file và đọc
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.pptx', '.ppt']:
            return PresentationReader.read_pptx(file_path)
        elif ext == '.pdf':
            return PresentationReader.read_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}. Supported: .pptx, .ppt, .pdf")
    
    @staticmethod
    def _convert_ppt_to_pdf(ppt_path: str, pdf_path: str):
        """Convert PPT/PPTX to PDF.
        - Windows: dùng COM Automation (yêu cầu Microsoft PowerPoint).
        - Linux/macOS: dùng LibreOffice headless.
        """
        ppt_path = os.path.abspath(ppt_path)
        pdf_path  = os.path.abspath(pdf_path)
        print(f"🔄 Converting PPT to PDF: {ppt_path} -> {pdf_path}")

        if platform.system() == "Windows":
            PresentationReader._convert_ppt_to_pdf_windows(ppt_path, pdf_path)
        else:
            PresentationReader._convert_ppt_to_pdf_libreoffice(ppt_path, pdf_path)

    @staticmethod
    def _convert_ppt_to_pdf_windows(ppt_path: str, pdf_path: str):
        """Windows-only: dùng COM Automation."""
        try:
            import pythoncom
            import comtypes.client
        except ImportError:
            raise Exception(
                "comtypes / pythoncom không được cài. "
                "Chạy: pip install comtypes pywin32"
            )

        pythoncom.CoInitialize()
        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
            presentation.SaveAs(pdf_path, 32)  # 32 = ppSaveAsPDF
            presentation.Close()
        except Exception as e:
            raise Exception(
                f"PPT Conversion failed (COM): {e}. "
                "Ensure Microsoft PowerPoint is installed."
            )
        finally:
            pythoncom.CoUninitialize()

    @staticmethod
    def _convert_ppt_to_pdf_libreoffice(ppt_path: str, pdf_path: str):
        """Linux/macOS: dùng LibreOffice headless."""
        output_dir = os.path.dirname(pdf_path)

        # Tìm LibreOffice executable
        candidates = ["libreoffice", "soffice", "/usr/bin/libreoffice", "/usr/bin/soffice"]
        lo_bin = None
        for c in candidates:
            try:
                result = subprocess.run(
                    [c, "--version"],
                    capture_output=True, timeout=5
                )
                if result.returncode == 0:
                    lo_bin = c
                    break
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue

        if lo_bin is None:
            raise Exception(
                "LibreOffice không được tìm thấy. "
                "Cài bằng: sudo apt-get install -y libreoffice"
            )

        cmd = [
            lo_bin,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            ppt_path
        ]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode != 0:
                raise Exception(
                    f"LibreOffice exited with code {result.returncode}:\n"
                    f"{result.stderr}"
                )
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timed out (>120s).")

        # LibreOffice đặt tên file theo tên gốc, đảm bảo đúng đường dẫn đầu ra
        generated = os.path.join(
            output_dir,
            os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf"
        )
        if generated != pdf_path and os.path.exists(generated):
            os.rename(generated, pdf_path)

        if not os.path.exists(pdf_path):
            raise Exception(f"Conversion succeeded but output file not found: {pdf_path}")

    @staticmethod
    def extract_slide_images(file_path: str, output_dir: str) -> List[str]:
        """
        Extract slides/pages as PNG images using pymupdf
        Handles PPT files by auto-converting to PDF first.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        os.makedirs(output_dir, exist_ok=True)
        
        ext = os.path.splitext(file_path)[1].lower()
        
        pdf_path = file_path
        
        # Handle PPT files
        if ext in ['.pptx', '.ppt']:
            pdf_path = os.path.splitext(file_path)[0] + ".pdf"
            
            should_convert = False
            if not os.path.exists(pdf_path):
                should_convert = True
            elif os.path.getmtime(file_path) > os.path.getmtime(pdf_path):
                should_convert = True
                
            if should_convert:
                PresentationReader._convert_ppt_to_pdf(file_path, pdf_path)
        
        # Extract slides from PDF using pymupdf
        try:
            # Suppress MuPDF stderr warnings (e.g. "No common ancestor in structure tree")
            # These are non-fatal PDF structure warnings that don't affect rendering
            fitz.TOOLS.mupdf_display_errors(False)
            doc = fitz.open(pdf_path)
            fitz.TOOLS.mupdf_display_errors(True)  # Restore for other operations
            image_paths = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Render page to image at high resolution
                zoom = 2
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)
                
                image_path = os.path.join(output_dir, f'slide_{page_num + 1}.png')
                pix.save(image_path)
                image_paths.append(image_path)
            
            doc.close()
            print(f"✅ Successfully extracted {len(image_paths)} slides from {pdf_path}")
            return image_paths
            
        except Exception as e:
            raise Exception(f"Error extracting slides: {str(e)}")
