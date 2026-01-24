"""
Module để đọc nội dung từ file PowerPoint và PDF
"""
from pptx import Presentation
from pypdf import PdfReader
import os
from typing import List, Dict

class PresentationReader:
    """Đọc nội dung từ file PPT và PDF"""
    
    @staticmethod
    def read_pptx(file_path: str) -> List[Dict]:
        """
        Đọc file PowerPoint và trả về danh sách slides
        
        Args:
            file_path: Đường dẫn đến file .pptx hoặc .ppt
        
        Returns:
            List[Dict]: [{
                'slide_num': int,
                'content': str,
                'notes': str (optional),
                'total_slides': int
            }]
        """
        try:
            prs = Presentation(file_path)
            slides_data = []
            
            for idx, slide in enumerate(prs.slides, 1):
                content_parts = []
                
                # Đọc text từ shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text.strip())
                
                # Đọc notes (nếu có)
                notes = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes = notes_slide.notes_text_frame.text
                
                content = "\n".join(content_parts)
                
                slides_data.append({
                    'slide_num': idx,
                    'content': content,
                    'notes': notes,
                    'total_slides': len(prs.slides)
                })
            
            return slides_data
        except Exception as e:
            raise Exception(f"Error reading PPTX: {str(e)}")
    
    @staticmethod
    def read_pdf(file_path: str) -> List[Dict]:
        """
        Đọc file PDF và trả về danh sách pages
        
        Args:
            file_path: Đường dẫn đến file .pdf
        
        Returns:
            List[Dict]: [{
                'slide_num': int,
                'content': str,
                'total_slides': int
            }]
        """
        try:
            reader = PdfReader(file_path)
            pages_data = []
            
            for idx, page in enumerate(reader.pages, 1):
                content = page.extract_text()
                
                pages_data.append({
                    'slide_num': idx,
                    'content': content.strip(),
                    'notes': "",  # PDF không có notes
                    'total_slides': len(reader.pages)
                })
            
            return pages_data
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
    
    @staticmethod
    def extract_text_from_file(file_path: str) -> List[Dict]:
        """
        Tự động detect loại file và đọc
        
        Args:
            file_path: Đường dẫn đến file
        
        Returns:
            List[Dict]: Danh sách slides/pages
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.pptx', '.ppt']:
            return PresentationReader.read_pptx(file_path)
        elif ext == '.pdf':
            return PresentationReader.read_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}. Supported: .pptx, .ppt, .pdf")

